import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pandas as pd
import json

def extract_text(file_path: str, filename: str) -> str:
    ext = filename.split(".")[-1].lower()

    if ext == "pdf":
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text

    elif ext == "docx":
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    elif ext == "txt" or ext == "md":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    elif ext == "csv":
        df = pd.read_csv(file_path)
        return df.to_string()

    elif ext == "xlsx":
        df = pd.read_excel(file_path)
        return df.to_string()

    elif ext == "json":
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
            return json.dumps(data, indent=2)

    elif ext == "pptx":
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    else:
        raise ValueError(f"Unsupported file type: {ext}")